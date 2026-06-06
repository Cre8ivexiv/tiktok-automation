from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Iterable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "outputs"
PPTX_PATH = OUT_DIR / "AIE_Autonomous_Security_Engineer_Pitch_Deck.pptx"
BRIEF_PATH = OUT_DIR / "AIE_Pitch_Deck_Brief.md"
SCRIPT_PATH = OUT_DIR / "AIE_3_Minute_Speaker_Script.md"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"
FONT_MONO = "Consolas"

BG = "07111F"
PANEL = "0D182A"
PANEL_2 = "111F36"
EDGE = "1C3554"
TEXT = "F5FAFF"
MUTED = "A8BEDA"
ACCENT = "31C5FF"
ACCENT_2 = "1E7DFF"
ACCENT_SOFT = "12385B"
SUCCESS = "22D3A6"
WARNING = "FF6B35"
WARNING_SOFT = "3E1A16"
GRID = "17314F"


@dataclass(frozen=True)
class SlideSpec:
    title: str
    content: list[str]
    visual: str
    notes: str
    animation: str
    sources: list[str] = field(default_factory=list)
    transition: str = "fade"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(BG)

    top_glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(8.8),
        Inches(-0.6),
        Inches(5.2),
        Inches(5.2),
    )
    style_shape(top_glow, fill=ACCENT_2, line=ACCENT_2, transparency=0.88)

    left_glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(-1.0),
        Inches(4.7),
        Inches(4.8),
        Inches(4.8),
    )
    style_shape(left_glow, fill=ACCENT, line=ACCENT, transparency=0.92)

    for offset in range(5):
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.45 + (offset * 0.1)),
            Inches(0.55 + (offset * 0.45)),
            Inches(12.0),
            Inches(0.01),
        )
        style_shape(line, fill=GRID, line=GRID, transparency=0.72)

    for idx in range(9):
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(11.55 + ((idx % 3) * 0.18)),
            Inches(6.25 + ((idx // 3) * 0.18)),
            Inches(0.05),
            Inches(0.05),
        )
        style_shape(dot, fill=ACCENT, line=ACCENT, transparency=0.2)


def style_shape(shape, *, fill: str, line: str, transparency: float = 0.0, line_width: float = 1.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    set_transparency(shape.fill, transparency)
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(line_width)


def set_transparency(fill, value: float) -> None:
    try:
        fill.transparency = value
    except Exception:
        pass


def add_panel(slide, x: float, y: float, w: float, h: float, *, fill: str = PANEL, line: str = EDGE, transparency: float = 0.0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    style_shape(shape, fill=fill, line=line, transparency=transparency, line_width=1.15)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 16,
    color_hex: str = TEXT,
    font_name: str = FONT_BODY,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    italic: bool = False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color(color_hex)
    return box


def add_paragraphs(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Iterable[str],
    *,
    size: float = 14,
    color_hex: str = TEXT,
    font_name: str = FONT_BODY,
    bold_first: bool = False,
    bullet: str | None = None,
    gap_after: float = 3,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.clear()
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.space_after = Pt(gap_after)
        run = p.add_run()
        run.text = f"{bullet} {line}" if bullet else line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold_first and idx == 0
        run.font.color.rgb = color(color_hex)
    return box


def add_section_tag(slide, label: str) -> None:
    tag = add_panel(slide, 0.72, 0.42, 2.35, 0.34, fill=ACCENT_SOFT, line=ACCENT, transparency=0.15)
    tag.line.width = Pt(1.0)
    add_text(
        slide,
        0.82,
        0.475,
        2.1,
        0.2,
        label.upper(),
        size=9,
        font_name=FONT_BODY,
        bold=True,
        color_hex=ACCENT,
    )


def add_title(slide, title: str, *, kicker: str | None = None) -> None:
    if kicker:
        add_section_tag(slide, kicker)
    add_text(slide, 0.72, 0.88, 7.5, 0.9, title, size=26, font_name=FONT_HEAD, bold=True)


def add_footer(slide, slide_no: int) -> None:
    add_text(slide, 0.72, 7.02, 3.2, 0.2, "AIE | BCU Enterprise Practice Project", size=8.5, color_hex=MUTED)
    add_text(slide, 12.3, 7.02, 0.35, 0.2, str(slide_no), size=8.5, color_hex=MUTED, align=PP_ALIGN.RIGHT)


def add_source(slide, text: str) -> None:
    add_text(slide, 0.72, 6.78, 11.6, 0.22, text, size=8.3, color_hex=MUTED)


def add_pill(slide, x: float, y: float, w: float, h: float, label: str, *, fill: str = PANEL_2, line: str = EDGE, text_color: str = TEXT):
    pill = add_panel(slide, x, y, w, h, fill=fill, line=line)
    pill.line.width = Pt(0.9)
    add_text(slide, x + 0.1, y + 0.07, w - 0.2, h - 0.1, label, size=10.5, color_hex=text_color, bold=True)
    return pill


def add_stat_card(slide, x: float, y: float, w: float, h: float, value: str, label: str, *, tone: str = ACCENT):
    add_panel(slide, x, y, w, h, fill=PANEL, line=tone)
    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.07),
        Inches(h),
    )
    style_shape(accent_bar, fill=tone, line=tone)
    add_text(slide, x + 0.22, y + 0.16, w - 0.35, 0.44, value, size=23, color_hex=tone, font_name=FONT_HEAD, bold=True)
    add_text(slide, x + 0.22, y + 0.66, w - 0.35, h - 0.72, label, size=11.2, color_hex=MUTED)


def add_terminal_card(slide, x: float, y: float, w: float, h: float, header: str, lines: list[str]) -> None:
    add_panel(slide, x, y, w, h, fill=PANEL, line=EDGE)
    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.34),
    )
    style_shape(top, fill=PANEL_2, line=PANEL_2)
    for idx, shade in enumerate(["FF5F56", "FFBD2E", "27C93F"]):
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x + 0.14 + (idx * 0.16)),
            Inches(y + 0.11),
            Inches(0.08),
            Inches(0.08),
        )
        style_shape(dot, fill=shade, line=shade)
    add_text(slide, x + 0.56, y + 0.07, w - 0.8, 0.18, header, size=9.5, color_hex=MUTED)
    add_paragraphs(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.66, lines, size=11.5, color_hex=TEXT, font_name=FONT_MONO, gap_after=6)


def add_warning_panel(slide, x: float, y: float, w: float, h: float) -> None:
    add_panel(slide, x, y, w, h, fill=PANEL, line=WARNING)
    tri = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        Inches(x + 0.35),
        Inches(y + 0.35),
        Inches(1.1),
        Inches(1.0),
    )
    style_shape(tri, fill=WARNING_SOFT, line=WARNING)
    add_text(slide, x + 0.74, y + 0.56, 0.25, 0.25, "!", size=18, color_hex=WARNING, font_name=FONT_HEAD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 1.7, y + 0.38, w - 2.05, 0.32, "Faster does not always mean safer.", size=16, font_name=FONT_HEAD, bold=True)
    add_paragraphs(
        slide,
        x + 1.7,
        y + 0.88,
        w - 2.05,
        h - 1.1,
        [
            "AI boosts code velocity and lowers friction.",
            "Trust, validation, and secure delivery are not keeping pace.",
            "Security debt is now created at generation speed.",
        ],
        size=11.3,
        color_hex=MUTED,
    )


def add_ring_badge(slide, x: float, y: float, d: float, label: str, sublabel: str | None = None) -> None:
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    style_shape(outer, fill=PANEL_2, line=ACCENT, transparency=0.15, line_width=2.0)
    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.13), Inches(y + 0.13), Inches(d - 0.26), Inches(d - 0.26))
    style_shape(inner, fill=PANEL, line=EDGE, transparency=0.05)
    add_text(slide, x + 0.25, y + 0.4, d - 0.5, 0.48, label, size=21, font_name=FONT_HEAD, bold=True, align=PP_ALIGN.CENTER)
    if sublabel:
        add_text(slide, x + 0.22, y + 0.92, d - 0.44, 0.24, sublabel, size=9.2, color_hex=MUTED, align=PP_ALIGN.CENTER)


def add_cloud_cluster(slide, x: float, y: float) -> None:
    for dx, dy, size in [(0.0, 0.28, 0.62), (0.34, 0.08, 0.8), (0.78, 0.24, 0.58)]:
        part = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x + dx),
            Inches(y + dy),
            Inches(size),
            Inches(size),
        )
        style_shape(part, fill=ACCENT_SOFT, line=ACCENT, transparency=0.15)
    base = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x + 0.14),
        Inches(y + 0.54),
        Inches(1.24),
        Inches(0.42),
    )
    style_shape(base, fill=ACCENT_SOFT, line=ACCENT, transparency=0.15)
    try:
        base.adjustments[0] = 0.2
    except Exception:
        pass


def add_code_lines(slide, x: float, y: float, widths: list[float], *, tone: str = ACCENT, mono_label: str | None = None) -> None:
    if mono_label:
        add_text(slide, x, y - 0.15, 1.4, 0.16, mono_label, size=8.5, color_hex=MUTED, font_name=FONT_MONO)
    for idx, width in enumerate(widths):
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y + (idx * 0.18)),
            Inches(width),
            Inches(0.05),
        )
        style_shape(line, fill=tone, line=tone, transparency=0.18)


def style_table_cell(cell, *, fill: str, text_value: str, size: float, bold: bool = False, font_name: str = FONT_BODY, color_hex: str = TEXT, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = color(fill)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text_value
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(color_hex)
    cell.text_frame.word_wrap = True
    cell.margin_left = Pt(5)
    cell.margin_right = Pt(5)
    cell.margin_top = Pt(4)
    cell.margin_bottom = Pt(4)


def add_donut_chart(slide, x: float, y: float, w: float, h: float, categories: list[str], values: list[int], colors: list[str]):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Series 1", values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.hole_size = 68
    plot.first_slice_angle = 270
    plot.has_data_labels = True
    plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.name = FONT_BODY
    plot.data_labels.font.color.rgb = color(TEXT)
    series = chart.series[0]
    for point, point_color in zip(series.points, colors, strict=True):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color(point_color)
        point.format.line.color.rgb = color(point_color)
    return chart


def add_bar_chart(slide, x: float, y: float, w: float, h: float, categories: list[str], values: list[int], point_colors: list[str], *, title: str | None = None):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Signal", values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False
    chart.category_axis.tick_labels.font.name = FONT_BODY
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = color(TEXT)
    chart.category_axis.format.line.color.rgb = color(GRID)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.major_unit = 20
    chart.value_axis.tick_labels.font.name = FONT_BODY
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.color.rgb = color(MUTED)
    chart.value_axis.major_gridlines.format.line.color.rgb = color(GRID)
    chart.value_axis.format.line.color.rgb = color(GRID)
    plot = chart.plots[0]
    plot.gap_width = 45
    plot.has_data_labels = True
    plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.color.rgb = color(TEXT)
    series = chart.series[0]
    for point, point_color in zip(series.points, point_colors, strict=True):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color(point_color)
        point.format.line.color.rgb = color(point_color)
    if title:
        add_text(slide, x, y - 0.26, w, 0.18, title, size=10, color_hex=MUTED)
    return chart


def add_line_chart(slide, x: float, y: float, w: float, h: float, categories: list[str], values: list[int]):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Spending", values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False
    chart.category_axis.tick_labels.font.name = FONT_BODY
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = color(TEXT)
    chart.category_axis.format.line.color.rgb = color(GRID)
    chart.value_axis.tick_labels.font.name = FONT_BODY
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.color.rgb = color(MUTED)
    chart.value_axis.major_gridlines.format.line.color.rgb = color(GRID)
    chart.value_axis.format.line.color.rgb = color(GRID)
    series = chart.series[0]
    series.format.line.color.rgb = color(ACCENT)
    series.format.line.width = Pt(2.5)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.ABOVE
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.color.rgb = color(TEXT)
    return chart


def add_step_box(slide, x: float, y: float, w: float, h: float, title: str, text: str, *, tone: str = ACCENT) -> None:
    add_panel(slide, x, y, w, h, fill=PANEL, line=tone)
    add_text(slide, x + 0.18, y + 0.14, w - 0.36, 0.24, title, size=12.5, font_name=FONT_HEAD, bold=True)
    add_text(slide, x + 0.18, y + 0.42, w - 0.36, h - 0.52, text, size=10.3, color_hex=MUTED)


def build_slide_1(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_section_tag(slide, "Strategic Framework Pitch")
    add_text(slide, 0.72, 1.0, 6.7, 0.95, "AIE: Autonomous Security Engineer", size=29, font_name=FONT_HEAD, bold=True)
    add_text(slide, 0.72, 1.84, 6.2, 0.56, "Security can’t stay manual. It has to be autonomous.", size=18.5, color_hex=ACCENT, font_name=FONT_HEAD, bold=True)
    add_paragraphs(slide, 0.72, 2.62, 4.2, 1.35, spec.content, size=12.2, color_hex=MUTED)
    add_pill(slide, 0.72, 4.62, 2.05, 0.42, "Artificial Intelligence Engineer", fill=ACCENT_SOFT, line=ACCENT, text_color=ACCENT)
    add_pill(slide, 2.92, 4.62, 1.82, 0.42, "Cloud-native security", fill=PANEL_2, line=EDGE)
    add_pill(slide, 4.88, 4.62, 1.76, 0.42, "Startup pitch deck", fill=PANEL_2, line=EDGE)

    hero = add_panel(slide, 8.0, 0.92, 4.55, 5.6, fill=PANEL, line=EDGE)
    hero.line.width = Pt(1.3)
    ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.24), Inches(1.62), Inches(2.1), Inches(2.1))
    style_shape(ring, fill=PANEL_2, line=ACCENT, transparency=0.12, line_width=2.0)
    ring_inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.48), Inches(1.86), Inches(1.62), Inches(1.62))
    style_shape(ring_inner, fill=BG, line=EDGE)
    add_text(slide, 9.68, 2.25, 1.24, 0.44, "AIE", size=23, font_name=FONT_HEAD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 9.44, 2.84, 1.78, 0.2, "autonomous cloud defender", size=8.6, color_hex=MUTED, align=PP_ALIGN.CENTER)
    add_cloud_cluster(slide, 8.72, 4.36)
    for cx, cy in [(8.92, 1.9), (11.18, 1.9), (8.92, 4.78), (11.18, 4.78)]:
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx), Inches(cy), Inches(0.18), Inches(0.18))
        style_shape(node, fill=ACCENT, line=ACCENT)
    for x1, y1, x2, y2 in [
        (9.02, 1.98, 9.95, 2.68),
        (11.18, 1.98, 10.55, 2.68),
        (9.01, 4.78, 9.95, 3.52),
        (11.16, 4.78, 10.55, 3.52),
    ]:
        connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = color(ACCENT)
        connector.line.width = Pt(1.35)
    add_code_lines(slide, 8.68, 5.48, [2.6, 1.8, 2.2, 1.45], tone=ACCENT_2, mono_label="continuous security loop")
    add_footer(slide, 1)


def build_slide_2(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Introduction and Hook")
    add_text(slide, 0.72, 1.64, 6.5, 0.4, "How many of you here code with AI?", size=24, font_name=FONT_HEAD, bold=True)
    add_text(slide, 0.72, 2.08, 6.5, 0.34, "AI helps developers build faster, but faster does not always mean safer.", size=12.8, color_hex=MUTED)
    add_terminal_card(
        slide,
        0.72,
        2.55,
        5.55,
        2.4,
        "developer-assistant session",
        [
            "> generate secure API login flow",
            "AI suggests endpoints, auth logic, and validation.",
            "Build speed rises instantly.",
            "Hidden security debt may ship just as fast.",
        ],
    )
    add_warning_panel(slide, 6.6, 2.55, 5.95, 2.4)
    add_stat_card(slide, 0.72, 5.28, 2.3, 1.18, "84%", "Developers use or plan to use AI tools", tone=ACCENT)
    add_stat_card(slide, 3.16, 5.28, 2.3, 1.18, "51%", "Professional developers use AI tools daily", tone=ACCENT_2)
    add_stat_card(slide, 5.6, 5.28, 2.3, 1.18, "46%", "Developers distrust AI output", tone=WARNING)
    add_pill(slide, 8.06, 5.28, 2.28, 0.52, "vibe coding is now mainstream", fill=ACCENT_SOFT, line=ACCENT, text_color=ACCENT)
    add_text(slide, 8.08, 5.92, 4.15, 0.44, "Adoption is accelerating ahead of trust and validation.", size=11.2, color_hex=MUTED)
    add_source(slide, "Source: Stack Overflow Developer Survey 2025 | survey.stackoverflow.co/2025/ai")
    add_footer(slide, 2)


def build_slide_3(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Problem")
    add_text(slide, 0.72, 1.58, 5.4, 0.56, "AI-generated and human-written code can both introduce hidden vulnerabilities.", size=18.8, font_name=FONT_HEAD, bold=True)
    add_pill(slide, 0.72, 2.3, 2.2, 0.45, "60% of breaches involve the human element", fill=WARNING_SOFT, line=WARNING, text_color=WARNING)
    add_pill(slide, 3.05, 2.3, 1.86, 0.45, "+34% vulnerability exploitation", fill=WARNING_SOFT, line=WARNING, text_color=WARNING)
    add_paragraphs(
        slide,
        0.72,
        3.0,
        4.45,
        2.6,
        [
            "Weak authentication",
            "SQL injection",
            "Exposed secrets",
            "Weak cryptography",
            "Cloud misconfiguration",
        ],
        size=11.8,
        color_hex=TEXT,
        bullet="•",
    )
    add_text(slide, 0.72, 5.88, 4.7, 0.42, "The issue is not just AI output. It is insecure delivery at cloud speed.", size=11.5, color_hex=MUTED)
    add_text(slide, 6.25, 1.82, 5.75, 0.26, "AI-generated code security", size=11, color_hex=MUTED)
    add_donut_chart(slide, 6.22, 2.05, 4.95, 3.95, ["Secure", "Failed security tests"], [55, 45], [ACCENT, WARNING])
    add_text(slide, 8.06, 3.7, 1.28, 0.4, "55%", size=23, font_name=FONT_HEAD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 7.76, 4.15, 1.92, 0.24, "roughly secure", size=9.2, color_hex=MUTED, align=PP_ALIGN.CENTER)
    add_source(slide, "Sources: Veracode 2025 GenAI Code Security Report; Verizon 2025 DBIR. 55% is inferred from Veracode’s 45% failure rate.")
    add_footer(slide, 3)


def build_slide_4(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Who Is Affected")
    cards = [
        ("DEV", "Developers", "AI-generated code and fast release cycles."),
        ("OPS", "DevOps", "Cloud access, deployment, and infrastructure risk."),
        ("SME", "Startups / SMEs", "Need enterprise-grade security without a large security team."),
        ("ENT", "Enterprise teams", "Complex systems, compliance pressure, broad attack surface."),
        ("SOC", "Security teams", "Alert fatigue, manual validation, remediation backlog."),
    ]
    x_positions = [0.72, 3.27, 5.82, 8.37, 10.92]
    for x, (badge, title, text) in zip(x_positions, cards, strict=True):
        add_panel(slide, x, 2.0, 2.1, 2.75, fill=PANEL, line=EDGE)
        add_ring_badge(slide, x + 0.56, 2.28, 0.96, badge)
        add_text(slide, x + 0.16, 3.46, 1.78, 0.28, title, size=12.5, font_name=FONT_HEAD, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.16, 3.86, 1.78, 0.68, text, size=9.6, color_hex=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, 0.72, 5.35, 12.0, 0.9, fill=PANEL_2, line=ACCENT)
    add_text(
        slide,
        1.0,
        5.63,
        11.35,
        0.28,
        "The people affected are not just security teams. It is every organisation shipping software faster than it can secure it.",
        size=14,
        font_name=FONT_HEAD,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 4)


def build_slide_5(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Why It Matters Now")
    add_text(slide, 0.72, 1.55, 5.8, 0.34, "Organisations are not just facing more attacks. They are facing faster attacks.", size=15.6, font_name=FONT_HEAD, bold=True)
    add_text(slide, 0.72, 2.02, 5.2, 0.22, "87% of intrusions now span multiple attack surfaces.", size=10.5, color_hex=ACCENT)
    add_stat_card(slide, 0.72, 2.55, 2.42, 1.18, "$4.88M", "Average global cost of a data breach in 2024", tone=ACCENT)
    add_stat_card(slide, 3.33, 2.55, 2.42, 1.18, "$10.5T", "Predicted annual global cybercrime cost in 2025", tone=ACCENT_2)
    add_stat_card(slide, 0.72, 4.05, 2.42, 1.18, "90%", "Breaches linked to preventable gaps or security misconfiguration", tone=WARNING)
    add_stat_card(slide, 3.33, 4.05, 2.42, 1.18, "4.76M", "Global cybersecurity workforce gap", tone=SUCCESS)
    add_bar_chart(
        slide,
        6.35,
        2.15,
        5.55,
        3.45,
        [
            "AI use / plan",
            "AI daily use",
            "Distrust AI output",
            "AI code failed tests",
            "Human element",
        ],
        [84, 51, 46, 45, 60],
        [ACCENT, ACCENT_2, WARNING, WARNING, SUCCESS],
        title="Core signals behind AIE",
    )
    add_text(slide, 6.35, 5.82, 5.65, 0.36, "Security pressure is increasing while teams remain understaffed. Manual checking does not scale.", size=11.3, color_hex=MUTED)
    add_source(slide, "Sources: IBM 2024; Cybersecurity Ventures 2025; Palo Alto Unit 42 2026; ISC2 2024; Stack Overflow 2025; Veracode 2025; Verizon 2025.")
    add_footer(slide, 5)


def build_slide_6(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Solution")
    add_text(slide, 0.72, 1.56, 5.6, 0.42, "AIE is an autonomous security platform that continuously monitors, learns, validates, and helps fix security risks across code and cloud environments.", size=12.4, color_hex=MUTED)
    center_x = 5.72
    center_y = 2.7
    add_ring_badge(slide, center_x, center_y, 1.7, "AIE", "autonomous security core")
    node_specs = [
        ("Detect", 4.15, 1.35),
        ("Learn", 7.7, 1.35),
        ("Simulate", 8.6, 3.34),
        ("Fix", 5.72, 4.62),
        ("Improve", 2.95, 3.34),
    ]
    for label, nx, ny in node_specs:
        add_ring_badge(slide, nx, ny, 1.18, label)
    for x1, y1, x2, y2 in [
        (5.25, 1.94, 6.08, 2.68),
        (7.7, 1.94, 7.15, 2.7),
        (8.75, 4.0, 7.65, 3.58),
        (6.3, 5.0, 6.55, 4.32),
        (4.04, 4.0, 5.25, 4.96),
        (4.13, 2.68, 5.25, 2.68),
    ]:
        connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = color(ACCENT)
        connector.line.width = Pt(1.5)
    side_x = 9.9
    add_step_box(slide, side_x, 1.5, 2.45, 0.88, "Secure code", "Detect flaws during development.", tone=ACCENT)
    add_step_box(slide, side_x, 2.52, 2.45, 0.88, "Secure cloud", "Surface posture gaps and risky permissions.", tone=ACCENT_2)
    add_step_box(slide, side_x, 3.54, 2.45, 0.88, "Learn from attackers", "Capture real-world behaviour with decoys.", tone=SUCCESS)
    add_step_box(slide, side_x, 4.56, 2.45, 0.88, "Validate defences", "Run controlled attack simulations.", tone=WARNING)
    add_footer(slide, 6)


def build_slide_7(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Product")
    quadrants = [
        ("1. Secure Coding Assistant", "Detects injection, authentication failures, secrets exposure, and weak cryptography during development.", ACCENT),
        ("2. Cloud Misconfiguration Scanner", "Monitors infrastructure for risky permissions, exposed services, and cloud posture gaps.", ACCENT_2),
        ("3. Honeypot Intelligence Layer", "Deploys decoys to capture and learn from attacker behaviour safely.", SUCCESS),
        ("4. Red vs Blue Simulation Engine", "Runs continuous attack validation to prove defences actually work.", WARNING),
    ]
    positions = [(0.72, 1.95), (6.7, 1.95), (0.72, 4.08), (6.7, 4.08)]
    for (title, body, tone), (x, y) in zip(quadrants, positions, strict=True):
        add_panel(slide, x, y, 5.7, 1.78, fill=PANEL, line=tone)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.2), Inches(y + 0.24), Inches(0.48), Inches(0.48))
        style_shape(badge, fill=tone, line=tone)
        add_text(slide, x + 0.86, y + 0.18, 4.45, 0.28, title, size=13.3, font_name=FONT_HEAD, bold=True)
        add_text(slide, x + 0.86, y + 0.54, 4.45, 0.64, body, size=10.6, color_hex=MUTED)
    add_pill(slide, 0.72, 6.18, 1.65, 0.38, "Injection", fill=WARNING_SOFT, line=WARNING, text_color=WARNING)
    add_pill(slide, 2.5, 6.18, 2.15, 0.38, "Cryptographic failures", fill=WARNING_SOFT, line=WARNING, text_color=WARNING)
    add_pill(slide, 4.8, 6.18, 2.1, 0.38, "Security misconfiguration", fill=WARNING_SOFT, line=WARNING, text_color=WARNING)
    add_pill(slide, 7.05, 6.18, 2.45, 0.38, "Identification/authentication", fill=WARNING_SOFT, line=WARNING, text_color=WARNING)
    add_source(slide, "Reference: OWASP Top 10 | owasp.org/Top10")
    add_footer(slide, 7)


def build_slide_8(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="How It Works")
    steps = [
        ("Detect", "Code and cloud scanned for vulnerabilities."),
        ("Learn", "Honeypots capture active exploit patterns."),
        ("Simulate", "Attack validation proves resilience."),
        ("Fix", "Remediation is suggested or applied."),
        ("Improve", "New intelligence updates future detection."),
    ]
    x = 0.72
    for idx, (title, body) in enumerate(steps):
        add_step_box(slide, x + (idx * 2.4), 2.15, 2.05, 1.2, title, body, tone=ACCENT if idx < 2 else ACCENT_2 if idx == 2 else SUCCESS if idx == 3 else WARNING)
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(x + 2.08 + (idx * 2.4)),
                Inches(2.56),
                Inches(0.24),
                Inches(0.36),
            )
            style_shape(arrow, fill=ACCENT, line=ACCENT)
    add_panel(slide, 0.72, 4.3, 12.0, 1.55, fill=PANEL_2, line=EDGE)
    add_text(slide, 0.95, 4.56, 2.2, 0.24, "Adaptive architecture", size=11.2, font_name=FONT_HEAD, bold=True)
    add_pill(slide, 3.0, 4.5, 1.72, 0.42, "Developer / AI code", fill=PANEL, line=ACCENT)
    add_pill(slide, 4.95, 4.5, 1.52, 0.42, "Code scanner", fill=PANEL, line=ACCENT)
    add_pill(slide, 6.65, 4.5, 1.58, 0.42, "Cloud scanner", fill=PANEL, line=ACCENT_2)
    add_pill(slide, 8.43, 4.5, 1.3, 0.42, "Honeypot", fill=PANEL, line=SUCCESS)
    add_pill(slide, 9.93, 4.5, 1.5, 0.42, "Simulation", fill=PANEL, line=WARNING)
    add_pill(slide, 11.63, 4.5, 0.86, 0.42, "Fix", fill=PANEL, line=SUCCESS)
    add_text(slide, 0.95, 5.12, 11.2, 0.34, "AIE becomes more useful over time because each layer teaches the others.", size=12.2, color_hex=TEXT)
    add_footer(slide, 8)


def build_slide_9(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Business Benefits")
    benefits = [
        ("Reduce Risk and Cost", "Catch flaws early and lower breach liability. AI and automation can reduce breach costs when used properly.", ACCENT),
        ("Operational Speed", "Reduce manual workload for security and DevOps, enabling faster and safer delivery.", ACCENT_2),
        ("Customer Trust", "Improve compliance readiness and protect reputation through stronger posture.", SUCCESS),
        ("Long-Term Resilience", "Turn real attack intelligence into stronger future protection.", WARNING),
    ]
    positions = [(0.72, 2.05), (6.62, 2.05), (0.72, 4.15), (6.62, 4.15)]
    for (title, body, tone), (x, y) in zip(benefits, positions, strict=True):
        add_panel(slide, x, y, 5.7, 1.65, fill=PANEL, line=tone)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.2), Inches(y + 0.22), Inches(0.44), Inches(0.44))
        style_shape(badge, fill=tone, line=tone)
        add_text(slide, x + 0.8, y + 0.16, 4.6, 0.28, title, size=13.1, font_name=FONT_HEAD, bold=True)
        add_text(slide, x + 0.8, y + 0.5, 4.6, 0.72, body, size=10.6, color_hex=MUTED)
    add_source(slide, "Source focus: IBM 2024 Cost of a Data Breach Report")
    add_footer(slide, 9)


def build_slide_10(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Competition and Advantage")
    add_text(slide, 0.72, 1.52, 11.4, 0.32, "Existing platforms focus on individual silos; AIE integrates them into one adaptive system.", size=14.8, font_name=FONT_HEAD, bold=True)
    table_shape = slide.shapes.add_table(6, 4, Inches(0.72), Inches(2.05), Inches(12.0), Inches(4.45))
    table = table_shape.table
    headers = ["Tool", "Main Focus", "Limitation", "AIE Advantage"]
    widths = [1.35, 2.45, 3.55, 4.65]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        style_table_cell(table.cell(0, col), fill=ACCENT_SOFT, text_value=header, size=10.8, bold=True, color_hex=ACCENT)
    rows = [
        ("Snyk", "Code security", "Strong developer focus, but mostly code-centric.", "Connects code signals to cloud, deception, and simulation."),
        ("Thinkst Canary", "Honeypots and deception", "Excellent decoys, but not a full remediation loop.", "Uses deception data to improve wider defences."),
        ("Prisma Cloud", "Cloud security / CNAPP", "Strong code-to-cloud posture, but not built around full attacker learning.", "Links posture management with adaptive attacker intelligence."),
        ("Tenable", "Exposure management", "Finds risk, but teams still need to act manually.", "Suggests and automates parts of the fix path."),
        ("AttackIQ", "Security validation", "Validates exposure, but not one loop across code, cloud, and decoys.", "Combines validation with continuous multi-layer learning."),
    ]
    for row_idx, row in enumerate(rows, start=1):
        fills = [PANEL, PANEL, PANEL, PANEL_2]
        colors = [TEXT, TEXT, MUTED, ACCENT]
        for col_idx, value in enumerate(row):
            style_table_cell(
                table.cell(row_idx, col_idx),
                fill=fills[col_idx],
                text_value=value,
                size=9.1,
                bold=col_idx == 0,
                color_hex=colors[col_idx],
            )
    add_source(slide, "Sources: snyk.io/platform | canary.tools | paloaltonetworks.com/prisma/cloud | tenable.com/products | attackiq.com")
    add_footer(slide, 10)


def build_slide_11(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Business Model")
    add_text(slide, 0.72, 1.52, 7.0, 0.36, "AIE is sold as SaaS, priced by team size, repositories, cloud assets, and infrastructure footprint.", size=12.4, color_hex=MUTED)
    tiers = [
        ("Starter", "£29 / month", ["Code scanning", "Basic dashboard", "Limited reports"], ACCENT),
        ("Growth", "£499 / month", ["Multi-repo support", "Cloud posture scans", "Honeypot integration"], ACCENT_2),
        ("Enterprise", "£2,500+ / month", ["Full automation", "Simulation engine", "Advanced support"], SUCCESS),
    ]
    x_positions = [0.72, 3.3, 5.88]
    for x, (name, price, bullets, tone) in zip(x_positions, tiers, strict=True):
        add_panel(slide, x, 2.2, 2.28, 3.65, fill=PANEL, line=tone)
        add_text(slide, x + 0.18, 2.42, 1.8, 0.24, name, size=13.4, font_name=FONT_HEAD, bold=True)
        add_text(slide, x + 0.18, 2.8, 1.8, 0.34, price, size=19, font_name=FONT_HEAD, bold=True, color_hex=tone)
        add_paragraphs(slide, x + 0.18, 3.42, 1.84, 1.4, bullets, size=10.3, color_hex=MUTED, bullet="•")
    add_panel(slide, 8.7, 2.0, 4.02, 3.95, fill=PANEL, line=EDGE)
    add_text(slide, 8.98, 2.22, 3.35, 0.24, "Worldwide information security spending", size=11.5, font_name=FONT_HEAD, bold=True)
    add_line_chart(slide, 9.02, 2.58, 3.3, 2.05, ["2024", "2025", "2026"], [193, 213, 240])
    add_text(slide, 9.02, 4.9, 3.25, 0.54, "Market timing supports entry: spending is projected to rise from $193B in 2024 to $240B in 2026.", size=10.6, color_hex=MUTED)
    add_source(slide, "Source: Gartner security spending forecast | 2024 $193B, 2025 $213B, 2026 $240B")
    add_footer(slide, 11)


def build_slide_12(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Investment Plan")
    add_stat_card(slide, 0.72, 2.0, 3.2, 1.35, "£100k – £150k", "Estimated year-one investment need", tone=ACCENT)
    add_text(slide, 0.72, 3.72, 2.9, 0.24, "Funding sources", size=12.5, font_name=FONT_HEAD, bold=True)
    add_paragraphs(
        slide,
        0.72,
        4.04,
        3.25,
        1.45,
        ["Founder contribution", "BCU innovation support", "Startup grants", "Seed investment", "Pilot partner support"],
        size=11,
        color_hex=MUTED,
        bullet="•",
    )
    add_panel(slide, 4.25, 1.9, 8.47, 4.95, fill=PANEL, line=EDGE)
    add_text(slide, 4.55, 2.15, 2.55, 0.24, "Cost allocation", size=12.5, font_name=FONT_HEAD, bold=True)
    add_donut_chart(
        slide,
        4.6,
        2.55,
        3.65,
        3.15,
        ["Product / research", "Cloud / testing", "Legal / compliance", "Marketing / sales", "Support / pilots"],
        [55, 18, 10, 12, 5],
        [ACCENT, ACCENT_2, SUCCESS, WARNING, "F9A26C"],
    )
    legend = [
        ("55% Product development and security research", ACCENT),
        ("18% Cloud hosting, logging, and testing", ACCENT_2),
        ("10% Legal, data protection, compliance", SUCCESS),
        ("12% Marketing and sales", WARNING),
        ("5% Support and pilot testing", "F9A26C"),
    ]
    for idx, (label, tone) in enumerate(legend):
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.6), Inches(2.5 + (idx * 0.56)), Inches(0.14), Inches(0.14))
        style_shape(badge, fill=tone, line=tone)
        add_text(slide, 8.82, 2.43 + (idx * 0.56), 3.25, 0.24, label, size=10.3, color_hex=MUTED)
    add_footer(slide, 12)


def build_slide_13(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Team and Fit")
    capability_lines = [
        "Cybersecurity degree learning and practical security modules",
        "Cloud security and shared-responsibility awareness",
        "Secure coding, vulnerability management, and OWASP understanding",
        "Security operations, incident response, and risk management interest",
        "Strong interest in automation, AI security, and developer tooling",
        "Built through the BCU Enterprise Practice Project journey",
    ]
    add_panel(slide, 0.72, 1.9, 5.5, 4.9, fill=PANEL, line=EDGE)
    add_text(slide, 1.0, 2.18, 3.6, 0.28, "Capability fit", size=13.5, font_name=FONT_HEAD, bold=True)
    for idx, line in enumerate(capability_lines):
        add_pill(slide, 1.0, 2.7 + (idx * 0.58), 4.85, 0.42, line, fill=PANEL_2, line=EDGE)
    add_panel(slide, 6.55, 1.9, 6.17, 4.9, fill=PANEL, line=EDGE)
    add_text(slide, 6.86, 2.18, 2.2, 0.28, "Roadmap to capability", size=13.5, font_name=FONT_HEAD, bold=True)
    roadmap = ["Learning", "MVP", "Pilot", "Adaptive Platform"]
    for idx, label in enumerate(roadmap):
        cx = 7.12 + (idx * 1.45)
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx), Inches(4.2), Inches(0.55), Inches(0.55))
        style_shape(node, fill=ACCENT if idx < 2 else SUCCESS if idx == 2 else WARNING, line=EDGE)
        add_text(slide, cx - 0.16, 4.84, 0.88, 0.2, label, size=10, align=PP_ALIGN.CENTER)
        if idx < len(roadmap) - 1:
            link = slide.shapes.add_connector(1, Inches(cx + 0.55), Inches(4.47), Inches(cx + 1.3), Inches(4.47))
            link.line.color.rgb = color(ACCENT)
            link.line.width = Pt(1.5)
    add_text(slide, 6.86, 2.82, 5.25, 0.72, "This project turns academic cybersecurity learning into a practical, industry-facing product idea with a clear technical and business direction.", size=11.8, color_hex=MUTED)
    add_footer(slide, 13)


def build_slide_14(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="Traction and Feedback")
    add_panel(slide, 0.72, 2.2, 3.8, 3.15, fill=PANEL, line=EDGE)
    add_text(slide, 1.02, 2.45, 1.6, 0.24, "Before", size=12.6, font_name=FONT_HEAD, bold=True)
    add_text(slide, 1.02, 2.82, 2.8, 0.28, "Simple scanner", size=18, font_name=FONT_HEAD, bold=True, color_hex=ACCENT)
    add_paragraphs(slide, 1.02, 3.4, 3.0, 1.1, ["One-off detection", "No adaptive learning", "Limited market differentiation"], size=10.8, color_hex=MUTED, bullet="•")
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(5.05), Inches(3.25), Inches(1.05), Inches(0.75))
    style_shape(arrow, fill=ACCENT, line=ACCENT)
    add_panel(slide, 6.45, 2.2, 5.9, 3.15, fill=PANEL, line=ACCENT)
    add_text(slide, 6.78, 2.45, 1.3, 0.24, "After", size=12.6, font_name=FONT_HEAD, bold=True)
    add_text(slide, 6.78, 2.82, 4.2, 0.28, "Autonomous adaptive platform", size=18, font_name=FONT_HEAD, bold=True, color_hex=ACCENT)
    add_paragraphs(
        slide,
        6.78,
        3.35,
        4.95,
        1.3,
        [
            "Secure coding + cloud remediation",
            "Honeypot intelligence",
            "Red vs blue validation loop",
            "Subscription SaaS model",
        ],
        size=10.8,
        color_hex=MUTED,
        bullet="•",
    )
    add_text(slide, 0.72, 5.95, 11.2, 0.34, "Feedback helped turn AIE from a tool into a platform.", size=14, font_name=FONT_HEAD, bold=True)
    add_footer(slide, 14)


def build_slide_15(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="The Ask")
    add_panel(slide, 0.72, 1.92, 4.75, 4.9, fill=PANEL, line=ACCENT)
    add_text(slide, 1.0, 2.18, 3.6, 0.28, "Support a 12-month pilot of AIE", size=17.2, font_name=FONT_HEAD, bold=True)
    add_paragraphs(
        slide,
        1.0,
        2.85,
        3.7,
        1.8,
        ["Pilot funding", "Cybersecurity mentor feedback", "Safe cloud test environment", "Pilot users from startups, SMEs, or university partners"],
        size=11,
        color_hex=MUTED,
        bullet="•",
    )
    add_text(slide, 1.0, 5.44, 3.3, 0.24, "Contact", size=11.3, font_name=FONT_HEAD, bold=True)
    add_text(slide, 1.0, 5.74, 3.8, 0.42, "Ibrahim Timilehin\nibrahim.timilehin@bcu.ac.uk", size=10.4, color_hex=MUTED)
    add_panel(slide, 5.8, 1.92, 6.92, 4.9, fill=PANEL, line=EDGE)
    add_text(slide, 6.08, 2.18, 2.35, 0.28, "What success looks like", size=13.5, font_name=FONT_HEAD, bold=True)
    success_items = [
        "Working MVP",
        "Validated code and cloud scanning workflow",
        "Honeypot intelligence prototype",
        "Clear pilot feedback and evidence for investment",
    ]
    for idx, item in enumerate(success_items):
        add_pill(slide, 6.08, 2.75 + (idx * 0.58), 5.8, 0.42, item, fill=PANEL_2, line=EDGE)
    add_text(slide, 6.08, 5.42, 5.6, 0.64, "Security can’t stay manual.\nIt has to be autonomous.", size=22, font_name=FONT_HEAD, bold=True, color_hex=ACCENT)
    add_footer(slide, 15)


def build_slide_16(slide, spec: SlideSpec) -> None:
    add_bg(slide)
    add_title(slide, spec.title, kicker="References")
    add_text(slide, 0.72, 1.55, 6.2, 0.32, "This pitch is supported by current industry research and cybersecurity reports.", size=12.4, color_hex=MUTED)
    left_refs = spec.content[:6]
    right_refs = spec.content[6:]
    add_panel(slide, 0.72, 2.02, 5.85, 4.75, fill=PANEL, line=EDGE)
    add_panel(slide, 6.87, 2.02, 5.85, 4.75, fill=PANEL, line=EDGE)
    add_paragraphs(slide, 0.98, 2.28, 5.2, 4.2, left_refs, size=9.5, color_hex=TEXT, gap_after=7)
    add_paragraphs(slide, 7.13, 2.28, 5.2, 4.2, right_refs, size=9.5, color_hex=TEXT, gap_after=7)
    add_text(slide, 0.72, 6.88, 10.7, 0.24, "These sources show why AIE is timely, relevant, and connected to real industry needs.", size=10.3, color_hex=MUTED)
    add_footer(slide, 16)


def slide_specs() -> list[SlideSpec]:
    return [
        SlideSpec(
            title="AIE: Autonomous Security Engineer",
            content=[
                "Ibrahim Timilehin",
                "BCU Enterprise Practice Project",
                "AIE = Artificial Intelligence Engineer",
                "Autonomous security for the vibe coding era",
            ],
            visual="Futuristic cloud-security hero panel with network nodes, AIE brand ring, and abstract cloud telemetry.",
            notes="Good afternoon, my name is Ibrahim Timilehin, and today I’m pitching AIE, which stands for Artificial Intelligence Engineer. AIE is an autonomous security engineer in the cloud, built for the way software is being developed today.",
            animation="Fade in title first, then subtitle, then name.",
            transition="fade",
        ),
        SlideSpec(
            title="The Era of Vibe Coding",
            content=[
                "Large question: How many of you here code with AI?",
                "Supporting line: AI helps developers build faster, but faster does not always mean safer.",
                "84% of developers use or plan to use AI tools.",
                "51% of professional developers use AI tools daily.",
                "46% of developers distrust AI output.",
            ],
            visual="Split-screen visual with an AI coding terminal on the left and a security warning panel on the right, supported by three stat cards.",
            notes="Before I explain the product, I want to ask a quick question. How many of you here have used AI when coding, what people now call vibe coding? Most of us either have, or we know someone who does. AI can generate code, explain errors, and speed up development, but adoption is moving faster than trust and security.",
            animation="Question appears first. Then show the three stats one by one.",
            sources=["https://survey.stackoverflow.co/2025/ai"],
            transition="wipe",
        ),
        SlideSpec(
            title="AI Helps Us Build Faster — Not Always Safer",
            content=[
                "AI-generated and human-written code can both introduce hidden vulnerabilities.",
                "Only 55% of AI-generated code was secure.",
                "Around 45% failed security tests.",
                "Around 60% of breaches still involve the human element.",
                "Vulnerability exploitation increased by 34%.",
            ],
            visual="Donut chart showing the 55 / 45 secure-versus-failed split, supported by risk labels for auth, SQL injection, secrets, cryptography, and cloud posture.",
            notes="Here is the problem. AI does not just help us code faster. It can also make us insecure faster. Veracode found that 45% of AI-generated code samples failed security tests, so only about 55% passed. But this is not only an AI problem. Developers who do not use AI still make mistakes too, such as weak authentication, exposed secrets, poor cryptography, and cloud misconfigurations.",
            animation="Animate the chart first, then reveal the risk icons.",
            sources=[
                "https://www.veracode.com/blog/genai-code-security-report/",
                "https://www.verizon.com/about/news/2025-data-breach-investigations-report",
            ],
            transition="fade",
        ),
        SlideSpec(
            title="Who Is Affected?",
            content=[
                "Developers using AI-generated code and fast development workflows",
                "DevOps teams responsible for deployment, cloud access, and infrastructure",
                "Cloud-first startups and SMEs without large security teams",
                "Enterprise software teams managing complex systems and compliance needs",
                "Security teams overloaded by alerts, vulnerabilities, and manual checks",
            ],
            visual="Five stakeholder cards spanning developer, DevOps, startup/SME, enterprise, and security teams.",
            notes="This problem affects more than just security teams. Developers are affected because they may ship vulnerable code without realising it. DevOps teams are affected because cloud infrastructure can be misconfigured. Startups and SMEs are affected because they often do not have large security teams. Enterprises are affected because their systems are complex and compliance matters. So really, it affects any organisation that is shipping software faster than it can secure it.",
            animation="Show each stakeholder group one by one.",
            transition="push",
        ),
        SlideSpec(
            title="Scaling Vulnerabilities at Speed",
            content=[
                "Quote: Organisations are not just facing more attacks. They are facing faster attacks.",
                "$4.88 million average global cost of a data breach in 2024.",
                "$10.5 trillion predicted annual global cybercrime cost by 2025.",
                "90% of breaches linked to misconfigurations or security gaps.",
                "87% of incidents span multiple attack surfaces.",
                "Global cybersecurity workforce gap: 4,763,963 people.",
            ],
            visual="Four metric cards plus a horizontal bar chart showing the core signals behind AIE.",
            notes="This matters because small flaws do not always stay small. IBM reported that the average data breach cost reached 4.88 million dollars in 2024. Cybercrime is predicted to cost the world 10.5 trillion dollars annually by 2025. Palo Alto links most breaches to misconfigurations or security gaps, and ISC2 estimates the cybersecurity workforce gap at over 4.7 million people. So organisations are under pressure to move fast, but many do not have enough security people to manually keep up.",
            animation="Metric cards appear one by one with a slight zoom, then the bar chart reveals.",
            sources=[
                "https://www.ibm.com/reports/data-breach",
                "https://cybersecurityventures.com/official-cybercrime-report-2025/",
                "https://www.paloaltonetworks.com/resources/research/unit-42-incident-response-report",
                "https://www.isc2.org/insights/2024/10/isc2-2024-cybersecurity-workforce-study",
            ],
            transition="fade",
        ),
        SlideSpec(
            title="AIE Closes the Gap",
            content=[
                "AIE is an autonomous security platform that continuously monitors, learns, validates, and helps fix security risks across code and cloud environments.",
                "Core loop: Detect → Learn → Simulate → Fix → Improve",
                "Focus areas: secure code, secure cloud, learn from attackers, validate defences",
            ],
            visual="Circular feedback loop with AIE at the centre and five connected loop nodes.",
            notes="That is the gap AIE is designed to close. AIE is not just another alert dashboard. It is an autonomous security layer that monitors code, checks cloud infrastructure, learns from real attacker behaviour, and validates whether defences work. The key idea is the loop: detect, learn, simulate, fix, and improve.",
            animation="Animate the loop clockwise.",
            transition="wipe",
        ),
        SlideSpec(
            title="One Platform, Four Core Layers",
            content=[
                "Secure Coding Assistant",
                "Cloud Misconfiguration Scanner",
                "Honeypot Intelligence Layer",
                "Red vs Blue Simulation Engine",
                "OWASP focus: injection, cryptographic failures, security misconfiguration, and identification/authentication failures",
            ],
            visual="Four-quadrant layout with distinct cards for code, cloud, honeypots, and simulation.",
            notes="AIE has four core layers. The first is a secure coding assistant that checks code as developers write it. The second is a cloud scanner that looks for risky permissions, exposed services, and posture gaps. The third is honeypot intelligence, where AIE deploys decoys to learn from attacker behaviour safely. The fourth is a simulation engine that runs red-team-versus-blue-team style validation, so organisations can test whether their defences actually work.",
            animation="Each quadrant fades in individually.",
            sources=["https://owasp.org/Top10/"],
            transition="fade",
        ),
        SlideSpec(
            title="Autonomous Adaptive Loop",
            content=[
                "Detect: code and cloud scanned for vulnerabilities.",
                "Learn: honeypots capture active exploit patterns.",
                "Simulate: attack validation proves resilience.",
                "Fix: remediation is suggested or applied.",
                "Improve: new intelligence updates future detection.",
            ],
            visual="Architecture flow from developer / AI code through scanners, honeypot, simulation, and remediation.",
            notes="The important part is how these features connect. AIE detects issues in code and cloud environments. It learns from honeypots. It runs simulations to test defences. Then it recommends or applies fixes. Over time, if one layer learns something new, that learning strengthens the other layers. That is what makes AIE adaptive rather than static.",
            animation="Use a flow animation from left to right.",
            transition="push",
        ),
        SlideSpec(
            title="Strategic Business Benefits",
            content=[
                "Reduce Risk and Cost",
                "Operational Speed",
                "Customer Trust",
                "Long-Term Resilience",
                "AI and automation can reduce breach costs when used properly.",
            ],
            visual="Four benefit cards with clean icon markers for risk, speed, trust, and resilience.",
            notes="The benefit to organisations is clear. AIE helps catch vulnerabilities earlier, reduce breach risk, reduce manual workload, and improve customer trust. This supports business goals because companies want to innovate quickly, but they cannot afford to let security become an afterthought. AIE helps security move at the speed of development.",
            animation="Benefit cards fade in one by one.",
            sources=["https://www.ibm.com/reports/data-breach"],
            transition="fade",
        ),
        SlideSpec(
            title="The Adaptive Advantage",
            content=[
                "Comparison across Snyk, Thinkst Canary, Prisma Cloud, Tenable, and AttackIQ",
                "Main point: AIE integrates code security, cloud posture, deception, and validation into one loop",
            ],
            visual="Modern competitor comparison table with the AIE advantage column highlighted.",
            notes="Existing tools are useful, but they mostly solve one part of the problem. Snyk focuses on code. Thinkst Canary focuses on deception. Prisma Cloud and Tenable focus on cloud and exposure management. AttackIQ focuses on validation. AIE’s advantage is not that these tools are useless. The advantage is that AIE brings the layers together, so learning in one area improves the others.",
            animation="Reveal competitor rows first, then highlight the AIE advantage column.",
            sources=[
                "https://snyk.io/platform/",
                "https://canary.tools/",
                "https://www.paloaltonetworks.com/prisma/cloud",
                "https://www.tenable.com/products",
                "https://www.attackiq.com/",
            ],
            transition="wipe",
        ),
        SlideSpec(
            title="SaaS Subscription Model",
            content=[
                "Starter — £29 / month",
                "Growth — £499 / month",
                "Enterprise — £2,500+ / month",
                "Security spending forecast: 2024 $193bn, 2025 $213bn, 2026 $240bn",
            ],
            visual="Three pricing cards beside a line chart for Gartner’s security spending forecast.",
            notes="AIE follows a subscription model. The price depends on team size and cloud footprint, which means small teams can start with basic scanning while larger organisations can pay for automation, simulations, and deeper integrations. This is also timely because Gartner projects security spending to continue rising from 193 billion dollars in 2024 to 240 billion dollars in 2026.",
            animation="Pricing cards slide up one by one. The line chart draws from 2024 to 2026.",
            sources=["https://www.gartner.com/en/newsroom/press-releases/2025-07-29-gartner-forecasts-worldwide-end-user-spending-on-information-security-to-total-213-billion-us-dollars-in-2025"],
            transition="fade",
        ),
        SlideSpec(
            title="Year One Investment Plan",
            content=[
                "Investment need: £100,000 — £150,000",
                "Funding sources: founder contribution, BCU innovation support, startup grants, seed investment, pilot partner support",
                "Cost allocation: 55% product/security research, 18% cloud/testing, 10% legal/compliance, 12% marketing/sales, 5% support/pilot testing",
            ],
            visual="Investment card with a donut chart and color-coded allocation legend.",
            notes="For year one, the estimated investment need is between 100,000 and 150,000 pounds. This would support product development, security research, cloud testing, legal compliance, marketing, and pilot support. The funding could come from a mix of founder contribution, BCU innovation support, startup grants, seed investment, and pilot partner support.",
            animation="Donut chart slices appear one by one.",
            transition="wipe",
        ),
        SlideSpec(
            title="Why I Can Build This",
            content=[
                "Cybersecurity background from degree and practical modules",
                "Knowledge of cloud security, misconfigurations, and shared responsibility",
                "Understanding of secure coding, vulnerabilities, and OWASP risks",
                "Experience and interest in security operations, incident response, and risk management",
                "Strong interest in automation, AI security, and developer-focused tools",
                "Built as part of my BCU Enterprise Practice Project and innovation journey",
            ],
            visual="Capability panel on the left and a learning-to-platform roadmap on the right.",
            notes="This project is not random for me. It connects directly to what I have been learning in cybersecurity, especially cloud security, secure coding, vulnerability management, security operations, and incident response. It also supports my employability because it shows that I can take technical learning and turn it into a practical business idea that solves a real industry problem.",
            animation="Skill icons fade in. Roadmap appears from left to right.",
            transition="push",
        ),
        SlideSpec(
            title="From Scanner to Autonomous Platform",
            content=[
                "The idea started as a security scanner.",
                "Feedback showed that a scanner alone would not be innovative enough.",
                "The idea improved by adding honeypots, cloud remediation, and red vs blue simulation.",
                "The business model shifted to subscription SaaS because security needs continuous protection.",
                "The product now focuses on continuous learning, not one-off detection.",
            ],
            visual="Before-and-after diagram showing the evolution from scanner to adaptive platform.",
            notes="The idea has developed through feedback. At first, it could have just been a scanner, but that would not be strong enough because many scanners already exist. The feedback helped me improve the idea by adding honeypots, cloud remediation, and simulation. That changed AIE from a one-off tool into a continuous security platform.",
            animation="Before side appears first, then arrow, then after side.",
            transition="wipe",
        ),
        SlideSpec(
            title="The ASK",
            content=[
                "Support a 12-month pilot of AIE with funding, mentorship, and access to a test environment.",
                "Need: pilot funding, cybersecurity mentor feedback, a safe cloud test environment, and pilot users.",
                "Success: working MVP, tested scanning workflow, honeypot prototype, pilot feedback, evidence for future investment.",
                "Closing line: Security can’t stay manual. It has to be autonomous.",
            ],
            visual="Bold final ask slide with support requirements, success criteria, and contact details.",
            notes="My ask is simple. I am asking for support for a 12-month pilot of AIE with funding, mentorship, and access to a safe test environment. The goal is to build a working MVP, test the core workflow, collect feedback, and prepare the project for future development. Security cannot stay manual. It has to be autonomous.",
            animation="Ask appears first, then the support needs, then the final quote.",
            transition="fade",
        ),
        SlideSpec(
            title="Data Sources",
            content=[
                "Stack Overflow Developer Survey 2025 | https://survey.stackoverflow.co/2025/ai",
                "Veracode GenAI Code Security Report | https://www.veracode.com/blog/genai-code-security-report/",
                "IBM Cost of a Data Breach Report 2024 / 2025 landing page | https://www.ibm.com/reports/data-breach",
                "IBM 2024 report summary | https://www.ibm.com/think/insights/whats-new-2024-cost-of-a-data-breach-report",
                "Verizon 2025 DBIR | https://www.verizon.com/about/news/2025-data-breach-investigations-report",
                "Gartner Security Spending Forecast | https://www.gartner.com/en/newsroom/press-releases/2025-07-29-gartner-forecasts-worldwide-end-user-spending-on-information-security-to-total-213-billion-us-dollars-in-2025",
                "ISC2 2024 Cybersecurity Workforce Study | https://www.isc2.org/insights/2024/10/isc2-2024-cybersecurity-workforce-study",
                "Palo Alto Networks Unit 42 2026 Incident Response Report | https://www.paloaltonetworks.com/resources/research/unit-42-incident-response-report",
                "OWASP Top 10 | https://owasp.org/Top10/",
                "Cybersecurity Ventures 2025 cybercrime report | https://cybersecurityventures.com/official-cybercrime-report-2025/",
                "Snyk | https://snyk.io/platform/ | Thinkst Canary | https://canary.tools/",
                "Prisma Cloud | https://www.paloaltonetworks.com/prisma/cloud | Tenable | https://www.tenable.com/products | AttackIQ | https://www.attackiq.com/",
            ],
            visual="Two-column clean references layout for presentation credibility and source traceability.",
            notes="This slide shows that the pitch is supported by recognised industry sources and not just opinion. The data comes from reports by Stack Overflow, Veracode, IBM, Verizon, Gartner, ISC2, Palo Alto, OWASP, Cybersecurity Ventures, and the relevant competitor platforms.",
            animation="Simple fade or no animation.",
            transition="fade",
        ),
    ]


BUILDERS: list[Callable] = [
    build_slide_1,
    build_slide_2,
    build_slide_3,
    build_slide_4,
    build_slide_5,
    build_slide_6,
    build_slide_7,
    build_slide_8,
    build_slide_9,
    build_slide_10,
    build_slide_11,
    build_slide_12,
    build_slide_13,
    build_slide_14,
    build_slide_15,
    build_slide_16,
]


SPEAKER_SCRIPT = """# AIE 3-Minute Speaker Script

## Slide 1
Good afternoon, my name is Ibrahim Timilehin, and today I’m pitching AIE, which stands for Artificial Intelligence Engineer. AIE is an autonomous security engineer in the cloud, built for the way software is being developed today.

## Slide 2
Before I explain the product, I want to ask a quick question. How many of you here have used AI when coding, what people now call vibe coding? Most of us either have, or we know someone who does. Stack Overflow’s 2025 survey shows that 84% of developers use or plan to use AI tools, 51% of professional developers use them daily, and 46% distrust the output.

## Slide 3
So here is the problem: AI helps teams code faster, but it can also make them insecure faster. Veracode found that 45% of AI-generated code samples failed security tests. At the same time, many breaches still involve the human element, and vulnerability exploitation keeps rising. The code may work, but it may still be unsafe.

## Slide 4
This affects more than security teams. It affects developers, DevOps teams, startups, SMEs, enterprise teams, and overloaded security operations teams. In short, it affects organisations shipping software faster than they can secure it.

## Slide 5
That matters now because the business impact is large. IBM reported an average breach cost of 4.88 million dollars in 2024. Cybercrime is forecast at 10.5 trillion dollars annually, and the cybersecurity workforce gap is still over 4.7 million people. The pressure to move fast is rising faster than the capacity to defend manually.

## Slide 6
That is the gap AIE is designed to close. AIE is not another dashboard. It is an autonomous security platform built around a simple loop: detect, learn, simulate, fix, and improve.

## Slide 7
The platform has four layers: a secure coding assistant, a cloud misconfiguration scanner, a honeypot intelligence layer, and a red-versus-blue simulation engine. Together, those layers cover the major security gaps modern teams face.

## Slide 8
The value comes from connection, not just coverage. AIE detects issues, learns from attacker behaviour, validates defences through simulation, and feeds that intelligence back into future detection and remediation. That makes the platform adaptive over time.

## Slide 9
For organisations, that means lower risk, less manual workload, better delivery speed, stronger customer trust, and better long-term resilience.

## Slide 10
Competitors already exist, but most operate in silos. Snyk focuses strongly on code, Thinkst Canary on deception, Prisma Cloud and Tenable on posture and exposure, and AttackIQ on validation. AIE’s advantage is bringing those layers into one adaptive loop.

## Slide 11
Commercially, AIE fits a SaaS model. Small teams can start with code scanning, growing teams can add cloud posture and honeypot features, and enterprise customers can adopt fuller automation and simulation. That fits a security market that Gartner expects to keep growing.

## Slide 12
For year one, the estimated investment need is between 100,000 and 150,000 pounds. That funding supports product development, security research, cloud testing, compliance, marketing, and pilot support.

## Slide 13
I believe I can build this because it directly connects to my cybersecurity learning, especially cloud security, secure coding, vulnerability management, security operations, and incident response. It turns academic learning into a practical industry solution.

## Slide 14
The idea has also improved through feedback. It started as a scanner, but feedback showed that was not innovative enough. Adding cloud remediation, honeypot intelligence, and simulation turned it into a stronger platform concept.

## Slide 15
So my ask is simple: support a 12-month pilot of AIE with funding, mentorship, and access to a safe cloud test environment. The goal is a working MVP, validated workflows, pilot feedback, and evidence for future investment.

## Slide 16
All of this is grounded in the industry sources listed on the final slide. Security can’t stay manual. It has to be autonomous.
"""


def write_brief(specs: list[SlideSpec]) -> None:
    lines = ["# AIE Pitch Deck Brief", ""]
    for idx, spec in enumerate(specs, start=1):
        lines.append(f"## Slide {idx} — {spec.title}")
        lines.append("")
        lines.append("**Slide title**")
        lines.append(spec.title)
        lines.append("")
        lines.append("**Slide content**")
        for item in spec.content:
            lines.append(f"- {item}")
        lines.append("")
        lines.append("**Suggested visual or chart**")
        lines.append(spec.visual)
        lines.append("")
        lines.append("**Speaker notes**")
        lines.append(spec.notes)
        lines.append("")
        lines.append("**Suggested animation**")
        lines.append(spec.animation)
        lines.append("")
        if spec.sources:
            lines.append("**Sources**")
            for source in spec.sources:
                lines.append(f"- {source}")
            lines.append("")
    BRIEF_PATH.write_text("\n".join(lines), encoding="utf-8")


def write_speaker_script() -> None:
    SCRIPT_PATH.write_text(SPEAKER_SCRIPT, encoding="utf-8")


def add_notes_and_transitions(pptx_path: Path, specs: list[SlideSpec]) -> None:
    import time

    from win32com.client import gencache

    transitions = {
        "fade": 1793,
        "wipe": 2817,
        "push": 3853,
    }
    placeholder_body = 2
    speed_medium = 2

    app = gencache.EnsureDispatch("PowerPoint.Application")
    app.Visible = 1
    presentation = app.Presentations.Open(str(pptx_path), WithWindow=True)
    try:
        time.sleep(1.0)
        for idx, spec in enumerate(specs, start=1):
            slide = None
            for _ in range(8):
                try:
                    slide = presentation.Slides(idx)
                    break
                except Exception:
                    time.sleep(0.35)
            if slide is None:
                raise RuntimeError(f"PowerPoint did not return slide {idx}")
            slide.SlideShowTransition.EntryEffect = transitions.get(spec.transition, 1793)
            slide.SlideShowTransition.Speed = speed_medium
            notes_shape = None
            for shape in slide.NotesPage.Shapes:
                try:
                    if shape.PlaceholderFormat.Type == placeholder_body:
                        notes_shape = shape
                        break
                except Exception:
                    continue
            if notes_shape is not None:
                notes_shape.TextFrame.TextRange.Text = spec.notes
        presentation.Save()
    finally:
        for closer in (presentation.Close, app.Quit):
            for _ in range(6):
                try:
                    closer()
                    break
                except Exception:
                    time.sleep(0.35)


def build_presentation(specs: list[SlideSpec]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "AIE — Autonomous Security Engineer in the Cloud"
    prs.core_properties.author = "Ibrahim Timilehin"
    prs.core_properties.subject = "BCU Enterprise Practice Project Pitch Deck"
    prs.core_properties.keywords = "AIE, cybersecurity, cloud security, AI, startup pitch"

    blank = prs.slide_layouts[6]
    for spec, builder in zip(specs, BUILDERS, strict=True):
        slide = prs.slides.add_slide(blank)
        builder(slide, spec)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    add_notes_and_transitions(PPTX_PATH, specs)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    specs = slide_specs()
    write_brief(specs)
    write_speaker_script()
    build_presentation(specs)
    print(f"Created: {PPTX_PATH}")
    print(f"Created: {BRIEF_PATH}")
    print(f"Created: {SCRIPT_PATH}")


if __name__ == "__main__":
    main()
